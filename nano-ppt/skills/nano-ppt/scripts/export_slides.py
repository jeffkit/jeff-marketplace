#!/usr/bin/env python3
"""
Export slide images to PPTX or PDF format
"""

import os
import sys
import argparse
import subprocess
from pathlib import Path
from typing import List, Optional

# Try importing required packages
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from PIL import Image
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False

try:
    import img2pdf
    IMG2PDF_AVAILABLE = True
except ImportError:
    IMG2PDF_AVAILABLE = False


def install_dependencies():
    """Install required Python dependencies"""
    packages = ['python-pptx', 'Pillow', 'img2pdf']
    for package in packages:
        print(f"📦 Installing {package}...")
        try:
            subprocess.check_call([sys.executable, "-m", "pip", "install", package, "-q"])
        except subprocess.CalledProcessError as e:
            print(f"❌ Failed to install {package}: {e}")
            return False
    return True


def get_slide_files(slides_dir: str) -> List[str]:
    """Get sorted list of slide image files"""
    slides_path = Path(slides_dir)
    if not slides_path.exists():
        raise FileNotFoundError(f"Slides directory not found: {slides_dir}")
    
    # Find all slide images (slide_01.png, slide_02.png, etc.)
    slide_files = sorted(
        [f for f in slides_path.glob("slide_*.png")],
        key=lambda x: int(x.stem.split("_")[1])
    )
    
    if not slide_files:
        raise FileNotFoundError(f"No slide images found in: {slides_dir}")
    
    return [str(f) for f in slide_files]


def export_to_pptx(slide_files: List[str], output_path: str, aspect_ratio: str = "16:9"):
    """Export slide images to PPTX"""
    if not PPTX_AVAILABLE:
        raise ImportError("python-pptx not available. Run with --check-deps first.")
    
    # Create presentation with specified aspect ratio
    prs = Presentation()
    
    # Set slide dimensions based on aspect ratio
    if aspect_ratio == "16:9":
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
    elif aspect_ratio == "4:3":
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
    else:
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
    
    # Use blank layout
    blank_layout = prs.slide_layouts[6]  # Blank slide
    
    for slide_file in slide_files:
        slide = prs.slides.add_slide(blank_layout)
        
        # Add image to fill the slide
        slide.shapes.add_picture(
            slide_file,
            Inches(0),
            Inches(0),
            width=prs.slide_width,
            height=prs.slide_height
        )
    
    prs.save(output_path)
    return output_path


def export_to_pdf(slide_files: List[str], output_path: str, aspect_ratio: str = "16:9"):
    """Export slide images to PDF with consistent page sizes"""
    if not IMG2PDF_AVAILABLE:
        raise ImportError("img2pdf not available. Run with --check-deps first.")
    
    if not PIL_AVAILABLE:
        raise ImportError("Pillow not available. Run with --check-deps first.")
    
    # Determine target size based on aspect ratio (at 150 DPI for good quality)
    # A4-ish landscape dimensions for 16:9
    if aspect_ratio == "16:9":
        target_width = 1920
        target_height = 1080
    else:  # 4:3
        target_width = 1600
        target_height = 1200
    
    # Normalize all images to same size and format
    temp_files = []
    for i, f in enumerate(slide_files):
        img = Image.open(f)
        
        # Convert to RGB if needed
        if img.mode in ('RGBA', 'P'):
            background = Image.new('RGB', img.size, (255, 255, 255))
            if img.mode == 'RGBA':
                background.paste(img, mask=img.split()[3])
            else:
                background.paste(img)
            img = background
        elif img.mode != 'RGB':
            img = img.convert('RGB')
        
        # Resize to target dimensions if different
        if img.size != (target_width, target_height):
            img = img.resize((target_width, target_height), Image.Resampling.LANCZOS)
        
        # Save as temporary file with consistent format
        temp_path = f.replace('.png', '_pdf_temp.jpg')
        img.save(temp_path, 'JPEG', quality=95, dpi=(150, 150))
        temp_files.append(temp_path)
    
    # Create PDF with explicit layout settings
    # Use img2pdf with consistent page size
    a4_width_pt = 841.89  # A4 landscape width in points
    a4_height_pt = 595.28  # A4 landscape height in points
    
    # Calculate layout to fit 16:9 in landscape A4-like page
    layout_fun = img2pdf.get_layout_fun(
        pagesize=(img2pdf.mm_to_pt(297), img2pdf.mm_to_pt(167))  # Custom 16:9 page
    )
    
    with open(output_path, "wb") as f:
        f.write(img2pdf.convert(temp_files, layout_fun=layout_fun))
    
    # Clean up temp files
    for f in temp_files:
        if os.path.exists(f):
            os.remove(f)
    
    return output_path


def main():
    parser = argparse.ArgumentParser(
        description="Export slide images to PPTX or PDF format"
    )
    parser.add_argument(
        "slides_dir",
        nargs='?',
        help="Directory containing slide images (slide_01.png, slide_02.png, ...)"
    )
    parser.add_argument(
        "output",
        nargs='?',
        help="Output file path (.pptx or .pdf)"
    )
    parser.add_argument(
        "--format",
        choices=["pptx", "pdf", "both"],
        default="pptx",
        help="Output format (default: pptx)"
    )
    parser.add_argument(
        "--aspect-ratio",
        default="16:9",
        choices=["16:9", "4:3"],
        help="Slide aspect ratio for PPTX (default: 16:9)"
    )
    parser.add_argument(
        "--check-deps",
        action="store_true",
        help="Check and install required dependencies"
    )
    
    args = parser.parse_args()
    
    # Handle dependency check
    if args.check_deps:
        print("🔍 Checking dependencies for export...")
        if install_dependencies():
            print("✅ All dependencies installed!")
        else:
            print("❌ Some dependencies failed to install")
            sys.exit(1)
        sys.exit(0)
    
    # Validate arguments
    if not args.slides_dir:
        parser.error("slides_dir is required")
    
    # Get slide files
    try:
        slide_files = get_slide_files(args.slides_dir)
        print(f"📊 Found {len(slide_files)} slides")
    except FileNotFoundError as e:
        print(f"❌ Error: {e}")
        sys.exit(1)
    
    # Determine output path
    slides_path = Path(args.slides_dir)
    parent_dir = slides_path.parent
    presentation_name = parent_dir.name
    
    # Export based on format
    try:
        if args.format in ["pptx", "both"]:
            output_pptx = args.output if args.output and args.output.endswith('.pptx') else str(parent_dir / f"{presentation_name}.pptx")
            export_to_pptx(slide_files, output_pptx, args.aspect_ratio)
            print(f"✅ PPTX exported: {output_pptx}")
        
        if args.format in ["pdf", "both"]:
            output_pdf = args.output if args.output and args.output.endswith('.pdf') else str(parent_dir / f"{presentation_name}.pdf")
            export_to_pdf(slide_files, output_pdf, args.aspect_ratio)
            print(f"✅ PDF exported: {output_pdf}")
            
    except ImportError as e:
        print(f"❌ Missing dependency: {e}")
        print("💡 Run with --check-deps to install required packages")
        sys.exit(1)
    except Exception as e:
        print(f"❌ Export failed: {e}")
        sys.exit(1)


if __name__ == "__main__":
    main()
